from flask import Flask, render_template, request, send_file, jsonify
import io
import os
import zipfile
import tempfile
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import arabic_reshaper
from bidi.algorithm import get_display
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pdf2docx import Converter
import pdfplumber
import pandas as pd
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from xhtml2pdf import pisa
import pytesseract
import google.generativeai as genai

# آپ کی چابی (API Key)
GEMINI_API_KEY = "AIzaSyBwnVZDJU57KIxsPg0020smcKV0nZGqO6I"

app = Flask(__name__)

try:
    font_path = "C:/Windows/Fonts/arial.ttf"
    if os.path.exists(font_path):
        pdfmetrics.registerFont(TTFont('UrduFont', font_path))
        DEFAULT_FONT = 'UrduFont'
    else:
        DEFAULT_FONT = 'Helvetica-Bold'
except:
    DEFAULT_FONT = 'Helvetica-Bold'

def process_urdu_text(text):
    try:
        reshaped_text = arabic_reshaper.reshape(text)
        bidi_text = get_display(reshaped_text)
        return bidi_text
    except:
        return text

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/tool')
def tool():
    return render_template('tool.html')

@app.route('/api/chat', methods=['POST'])
def chat_api():
    if not GEMINI_API_KEY: return jsonify({"error": "API Key is missing!"}), 400
    try:
        file = request.files.get('file'); question = request.form.get('question')
        if not file or not question: return jsonify({"error": "Please provide a file and a question."}), 400
        pdf_document = fitz.open(stream=file.read(), filetype="pdf"); extracted_text = ""
        for page_num in range(len(pdf_document)): extracted_text += pdf_document.load_page(page_num).get_text()
        pdf_document.close()
        if not extracted_text.strip(): return jsonify({"error": "No text found."}), 400
        extracted_text = extracted_text[:30000]; genai.configure(api_key=GEMINI_API_KEY)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        model = genai.GenerativeModel(available_models[0] if available_models else 'gemini-1.5-pro-latest')
        prompt = f"Document Text:\n\n{extracted_text}\n\nUser Question: {question}\n\nPlease answer based ONLY on the provided document."
        response = model.generate_content(prompt); return jsonify({"answer": response.text})
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/process', methods=['POST'])
def process_files():
    tool_name = request.form.get('tool_name', '').strip()
    files = request.files.getlist('files')

    if not files or files[0].filename == '':
        return "Please select a file!", 400

    if tool_name == "Merge PDF":
        merger = PdfMerger()
        for file in files: merger.append(file)
        output_pdf = io.BytesIO(); merger.write(output_pdf); merger.close(); output_pdf.seek(0)
        return send_file(output_pdf, as_attachment=True, download_name="Al_Musfillah_Merged.pdf", mimetype="application/pdf")

    elif tool_name == "Split PDF":
        uploaded_file = files[0]; reader = PdfReader(uploaded_file); memory_zip = io.BytesIO()
        with zipfile.ZipFile(memory_zip, 'w') as zf:
            for page_num in range(len(reader.pages)):
                writer = PdfWriter(); writer.add_page(reader.pages[page_num]); page_memory = io.BytesIO(); writer.write(page_memory); zf.writestr(f"Page_{page_num + 1}.pdf", page_memory.getvalue())
        memory_zip.seek(0); return send_file(memory_zip, as_attachment=True, download_name=f"Split_{uploaded_file.filename}.zip", mimetype="application/zip")

    elif tool_name == "Rotate PDF":
        uploaded_file = files[0]; rotation_angle = int(request.form.get('rotation_angle', 90)); reader = PdfReader(uploaded_file); writer = PdfWriter()
        for page in reader.pages: page.rotate(rotation_angle); writer.add_page(page)
        output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0); return send_file(output_pdf, as_attachment=True, download_name=f"Rotated_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "Delete PDF Pages":
        uploaded_file = files[0]; pages_string = request.form.get('pages_to_delete', '')
        try: pages_to_remove = [int(p.strip()) for p in pages_string.split(',') if p.strip().isdigit()]
        except: return "Invalid page numbers!", 400
        reader = PdfReader(uploaded_file); writer = PdfWriter()
        for index, page in enumerate(reader.pages):
            if (index + 1) not in pages_to_remove: writer.add_page(page)
        output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0); return send_file(output_pdf, as_attachment=True, download_name=f"Updated_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "Edit PDF":
        uploaded_file = files[0]; custom_text = request.form.get('custom_text', 'Al Musfillah PDF'); final_text = process_urdu_text(custom_text)
        try: x_pos = float(request.form.get('x_pos', 200)); y_pos = float(request.form.get('y_pos', 750))
        except: x_pos, y_pos = 200, 750
        packet = io.BytesIO(); can = canvas.Canvas(packet, pagesize=letter); can.setFont(DEFAULT_FONT, 24); can.setFillColorRGB(0.1, 0.3, 0.8); can.drawString(x_pos, y_pos, final_text); can.save(); packet.seek(0)
        new_pdf = PdfReader(packet); existing_pdf = PdfReader(uploaded_file); writer = PdfWriter()
        for i, page in enumerate(existing_pdf.pages):
            if i == 0: page.merge_page(new_pdf.pages[0])
            writer.add_page(page)
        output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0); return send_file(output_pdf, as_attachment=True, download_name=f"Edited_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "PDF Reader":
        uploaded_file = files[0]; file_data = io.BytesIO(uploaded_file.read()); file_data.seek(0); return send_file(file_data, as_attachment=False, download_name=uploaded_file.filename, mimetype="application/pdf")

    elif tool_name == "Number Pages":
        uploaded_file = files[0]; position = request.form.get('number_position', 'bottom_center'); existing_pdf = PdfReader(uploaded_file); writer = PdfWriter(); total_pages = len(existing_pdf.pages)
        for i, page in enumerate(existing_pdf.pages):
            packet = io.BytesIO(); page_width = float(page.mediabox.width); page_height = float(page.mediabox.height)
            can = canvas.Canvas(packet, pagesize=(page_width, page_height)); can.setFont(DEFAULT_FONT, 12); can.setFillColorRGB(0, 0, 0)
            text = f"Page {i + 1} of {total_pages}"; y_pos = page_height - 40 if 'top' in position else 30
            if 'left' in position: can.drawString(40, y_pos, text)
            elif 'right' in position: can.drawRightString(page_width - 40, y_pos, text)
            else: can.drawCentredString(page_width / 2, y_pos, text)
            can.save(); packet.seek(0); page.merge_page(PdfReader(packet).pages[0]); writer.add_page(page)
        output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0); return send_file(output_pdf, as_attachment=True, download_name=f"Numbered_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "Watermark PDF":
        uploaded_file = files[0]; watermark_text = request.form.get('watermark_text', 'Al Musfillah PDF'); final_watermark = process_urdu_text(watermark_text)
        try: opacity = float(request.form.get('watermark_opacity', 0.3))
        except: opacity = 0.3
        existing_pdf = PdfReader(uploaded_file); writer = PdfWriter()
        for page in existing_pdf.pages:
            packet = io.BytesIO(); page_width = float(page.mediabox.width); page_height = float(page.mediabox.height)
            can = canvas.Canvas(packet, pagesize=(page_width, page_height)); can.setFont(DEFAULT_FONT, 60); can.setFillAlpha(opacity); can.translate(page_width / 2, page_height / 2); can.rotate(45); can.drawCentredString(0, 0, final_watermark); can.save(); packet.seek(0)
            watermark_page = PdfReader(packet).pages[0]; page.merge_page(watermark_page); writer.add_page(page)
        output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0); return send_file(output_pdf, as_attachment=True, download_name=f"Watermarked_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "PDF to Word":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf: uploaded_file.save(temp_pdf.name); temp_pdf_path = temp_pdf.name
        temp_docx_path = temp_pdf_path.replace('.pdf', '.docx')
        try:
            cv = Converter(temp_pdf_path); cv.convert(temp_docx_path); cv.close(); return_data = io.BytesIO()
            with open(temp_docx_path, 'rb') as f: return_data.write(f.read())
            return_data.seek(0); os.remove(temp_pdf_path); os.remove(temp_docx_path)
            return send_file(return_data, as_attachment=True, download_name=f"{original_name}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        except Exception as e: return f"Error: {str(e)}", 500

    elif tool_name == "Protect PDF":
        uploaded_file = files[0]; password = request.form.get('pdf_password', '')
        if not password: return "Error: Password required!", 400
        reader = PdfReader(uploaded_file); writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        writer.encrypt(password); output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0)
        return send_file(output_pdf, as_attachment=True, download_name=f"Protected_{uploaded_file.filename}", mimetype="application/pdf")

    elif tool_name == "PDF to Excel":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf: uploaded_file.save(temp_pdf.name); temp_pdf_path = temp_pdf.name
        temp_excel_path = temp_pdf_path.replace('.pdf', '.xlsx')
        try:
            all_tables = []
            with pdfplumber.open(temp_pdf_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables: df = pd.DataFrame(table[1:], columns=table[0]); all_tables.append(df)
            if not all_tables: os.remove(temp_pdf_path); return "No tables found!", 400
            with pd.ExcelWriter(temp_excel_path, engine='openpyxl') as writer:
                for i, df in enumerate(all_tables): df.to_excel(writer, sheet_name=f'Table_{i+1}', index=False)
            return_data = io.BytesIO()
            with open(temp_excel_path, 'rb') as f: return_data.write(f.read())
            return_data.seek(0); os.remove(temp_pdf_path); os.remove(temp_excel_path)
            return send_file(return_data, as_attachment=True, download_name=f"{original_name}.xlsx", mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        except Exception as e: return f"Error: {str(e)}", 500

    elif tool_name == "PDF to PPT":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf: uploaded_file.save(temp_pdf.name); temp_pdf_path = temp_pdf.name
        temp_ppt_path = temp_pdf_path.replace('.pdf', '.pptx')
        try:
            prs = Presentation(); blank_slide_layout = prs.slide_layouts[6]; pdf_document = fitz.open(temp_pdf_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num); pix = page.get_pixmap(dpi=150); img_path = f"{temp_pdf_path}_{page_num}.png"; pix.save(img_path); slide = prs.slides.add_slide(blank_slide_layout); slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height); os.remove(img_path)
            pdf_document.close(); prs.save(temp_ppt_path); return_data = io.BytesIO()
            with open(temp_ppt_path, 'rb') as f: return_data.write(f.read())
            return_data.seek(0); os.remove(temp_pdf_path); os.remove(temp_ppt_path)
            return send_file(return_data, as_attachment=True, download_name=f"{original_name}.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e: return f"Error: {str(e)}", 500

    elif tool_name == "PDF to JPG":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf: uploaded_file.save(temp_pdf.name); temp_pdf_path = temp_pdf.name
        try:
            pdf_document = fitz.open(temp_pdf_path); memory_zip = io.BytesIO()
            with zipfile.ZipFile(memory_zip, 'w') as zf:
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num); pix = page.get_pixmap(dpi=300); img_bytes = pix.tobytes("jpeg"); zf.writestr(f"{original_name}_{page_num+1}.jpg", img_bytes)
            pdf_document.close(); os.remove(temp_pdf_path); memory_zip.seek(0)
            return send_file(memory_zip, as_attachment=True, download_name=f"{original_name}_JPG.zip", mimetype="application/zip")
        except Exception as e: return f"Error: {str(e)}", 500

    elif tool_name == "Word to PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_docx: uploaded_file.save(temp_docx.name); temp_docx_path = temp_docx.name
        temp_pdf_path = temp_docx_path.replace('.docx', '.pdf')
        try:
            convert(temp_docx_path, temp_pdf_path); return_data = io.BytesIO(); 
            with open(temp_pdf_path, 'rb') as f: return_data.write(f.read())
            return_data.seek(0); os.remove(temp_docx_path); 
            if os.path.exists(temp_pdf_path): os.remove(temp_pdf_path)
            return send_file(return_data, as_attachment=True, download_name=f"{original_name}.pdf", mimetype="application/pdf")
        except Exception as e: return f"Error: {str(e)}", 500

    elif tool_name == "JPG to PDF":
        try:
            images = []; 
            for file in files: img = Image.open(file).convert('RGB'); images.append(img)
            if not images: return "Error: No valid images selected!", 400
            output_pdf = io.BytesIO(); images[0].save(output_pdf, format="PDF", save_all=True, append_images=images[1:]); output_pdf.seek(0); original_name = files[0].filename.rsplit('.', 1)[0]
            return send_file(output_pdf, as_attachment=True, download_name=f"{original_name}_converted.pdf", mimetype="application/pdf")
        except Exception as e: return f"Error converting images to PDF.", 500

    elif tool_name == "HTML to PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        try:
            html_content = uploaded_file.read().decode('utf-8', errors='ignore'); output_pdf = io.BytesIO(); pisa_status = pisa.CreatePDF(html_content, dest=output_pdf)
            if pisa_status.err: return "Error: Could not process HTML.", 500
            output_pdf.seek(0)
            return send_file(output_pdf, as_attachment=True, download_name=f"{original_name}.pdf", mimetype="application/pdf")
        except Exception as e: return f"Error converting HTML to PDF.", 500

    elif tool_name in ["PDF to OCR", "PDF OCR"]:
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf: uploaded_file.save(temp_pdf.name); temp_pdf_path = temp_pdf.name
        temp_txt_path = temp_pdf_path.replace('.pdf', '.txt')
        try:
            pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            pdf_document = fitz.open(temp_pdf_path); extracted_text = ""
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num); pix = page.get_pixmap(dpi=300); mode = "RGBA" if pix.alpha else "RGB"; img = Image.frombytes(mode, [pix.width, pix.height], pix.samples); text = pytesseract.image_to_string(img); extracted_text += f"--- Page {page_num + 1} ---\n{text}\n\n"
            pdf_document.close()
            with open(temp_txt_path, 'w', encoding='utf-8') as f: f.write(extracted_text)
            return_data = io.BytesIO()
            with open(temp_txt_path, 'rb') as f: return_data.write(f.read())
            return_data.seek(0); os.remove(temp_pdf_path); os.remove(temp_txt_path)
            return send_file(return_data, as_attachment=True, download_name=f"{original_name}_OCR.txt", mimetype="text/plain")
        except Exception as e: return f"Error processing OCR: {str(e)}", 500

    elif tool_name == "Translate PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]; target_lang = request.form.get('target_language', 'Urdu')
        if not GEMINI_API_KEY: return "Error: API Key is missing!", 400
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf"); extracted_text = ""
            for page_num in range(min(10, len(pdf_document))): extracted_text += pdf_document.load_page(page_num).get_text()
            pdf_document.close()
            if not extracted_text.strip(): return "Error: No text found.", 400
            extracted_text = extracted_text[:15000]; genai.configure(api_key=GEMINI_API_KEY)
            available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
            model = genai.GenerativeModel(available_models[0] if available_models else 'gemini-1.5-pro-latest')
            prompt = f"Translate the following document text into {target_lang}. Return ONLY the translated text.\n\n{extracted_text}"
            response = model.generate_content(prompt); return_data = io.BytesIO()
            return_data.write('\ufeff'.encode('utf8')); return_data.write(response.text.encode('utf-8')); return_data.seek(0)
            return send_file(return_data, as_attachment=True, download_name=f"Translated_{target_lang}_{original_name}.txt", mimetype="text/plain")
        except Exception as e: return f"Translation Error: {str(e)}", 500

    elif tool_name == "Summarize PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]
        if not GEMINI_API_KEY: return "Error: API Key is missing!", 400
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf"); extracted_text = ""
            for page_num in range(min(15, len(pdf_document))): extracted_text += pdf_document.load_page(page_num).get_text()
            pdf_document.close()
            if not extracted_text.strip(): return "Error: No text found.", 400
            extracted_text = extracted_text[:25000]; genai.configure(api_key=GEMINI_API_KEY)
            available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
            model = genai.GenerativeModel(available_models[0] if available_models else 'gemini-1.5-pro-latest')
            prompt = f"Provide a comprehensive but concise summary in bullet points:\n\n{extracted_text}"
            response = model.generate_content(prompt); return_data = io.BytesIO()
            return_data.write('\ufeff'.encode('utf8')); return_data.write(response.text.encode('utf-8')); return_data.seek(0)
            return send_file(return_data, as_attachment=True, download_name=f"Summary_{original_name}.txt", mimetype="text/plain")
        except Exception as e: return f"Summarization Error: {str(e)}", 500

    elif tool_name == "Sign PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]; signature_file = request.files.get('signature_image')
        if not signature_file or signature_file.filename == '': return "Error: Upload a signature image!", 400
        try: page_num = int(request.form.get('sign_page', 1)) - 1; x_percent = float(request.form.get('sign_x_percent', 10)); y_percent = float(request.form.get('sign_y_percent', 80)); sign_width = float(request.form.get('sign_width', 150)); sign_height = float(request.form.get('sign_height', 50))
        except ValueError: page_num, x_percent, y_percent, sign_width, sign_height = 0, 10, 80, 150, 50
        try:
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            if page_num < 0 or page_num >= len(pdf_document): page_num = 0
            page = pdf_document.load_page(page_num); page_width = page.rect.width; page_height = page.rect.height
            x_pos = (x_percent / 100) * page_width; y_pos = (y_percent / 100) * page_height
            sig_bytes = signature_file.read(); rect = fitz.Rect(x_pos, y_pos, x_pos + sign_width, y_pos + sign_height)
            page.insert_image(rect, stream=sig_bytes); output_pdf = io.BytesIO(); pdf_document.save(output_pdf); pdf_document.close(); output_pdf.seek(0)
            return send_file(output_pdf, as_attachment=True, download_name=f"Signed_{original_name}.pdf", mimetype="application/pdf")
        except Exception as e: return f"Sign PDF Error: {str(e)}", 500

    elif tool_name == "Unlock PDF":
        uploaded_file = files[0]; original_name = uploaded_file.filename.rsplit('.', 1)[0]; password = request.form.get('pdf_password', '')
        if not password: return "Error: Please enter password!", 400
        try:
            reader = PdfReader(uploaded_file)
            if not reader.is_encrypted: return "Error: PDF is not locked.", 400
            if reader.decrypt(password) == 0: return "Error: Incorrect password!", 400
            writer = PdfWriter(); 
            for page in reader.pages: writer.add_page(page)
            output_pdf = io.BytesIO(); writer.write(output_pdf); output_pdf.seek(0)
            return send_file(output_pdf, as_attachment=True, download_name=f"Unlocked_{original_name}.pdf", mimetype="application/pdf")
        except Exception as e: return f"Unlock PDF Error: {str(e)}", 500

    # ==========================================
    # === 23. Redact PDF (نیا ٹول) ===
    # ==========================================
    elif tool_name == "Redact PDF":
        uploaded_file = files[0]
        original_name = uploaded_file.filename.rsplit('.', 1)[0]
        text_to_redact = request.form.get('redact_text', '').strip()

        if not text_to_redact:
            return "Error: Please enter the text you want to redact/hide!", 400

        try:
            # پی ڈی ایف کو کھولنا
            pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            text_found = False

            # تمام صفحات میں باری باری چیک کرنا
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                
                # مطلوبہ لفظ یا لائن تلاش کرنا
                areas = page.search_for(text_to_redact)
                
                if areas:
                    text_found = True
                    for rect in areas:
                        # اس لفظ پر کالا ڈبہ (Redaction Annotation) بنانا
                        page.add_redact_annot(rect, fill=(0, 0, 0))
                    
                    # یہ لائن اصلی کام کرتی ہے: کالا ڈبہ پکا کرتی ہے اور نیچے سے اصلی ٹیکسٹ اڑا دیتی ہے
                    page.apply_redactions()

            if not text_found:
                return f"Error: The exact text '{text_to_redact}' was not found in this PDF.", 400

            # فائل کو سیو کر کے بھیجنا
            output_pdf = io.BytesIO()
            pdf_document.save(output_pdf)
            pdf_document.close()
            output_pdf.seek(0)

            return send_file(output_pdf, as_attachment=True, download_name=f"Redacted_{original_name}.pdf", mimetype="application/pdf")
            
        except Exception as e:
            return f"Redact PDF Error: {str(e)}", 500


    return f"Error: Backend for '{tool_name}' is not ready yet!"

if __name__ == '__main__':
    app.run(debug=True)
